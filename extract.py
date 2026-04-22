#!/usr/bin/env python3
"""Extract text + images from .pptx files into markdown."""
import sys
import hashlib
from pathlib import Path
from collections import Counter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC_DIR = Path("/Users/rosta.nandor/Documents/jogalap")
PPT_DIR = SRC_DIR / "ppt"
DOCS_DIR = SRC_DIR / "docs"
IMG_DIR = DOCS_DIR / "images"

MAPPING = {
    "bevezetes_a_jog_szerepe_az_informacios_tarsadalomban20260220.pptx": "01-bevezetes",
    "Alapjogok_Jogi_alapismeretek_VIK_2026_tavasz.pptx": "02-alapjogok",
    "PT - VIKAdatvédelem26.pptx": "03-adatvedelem",
    "üzleti-működés-formái_jogi-alapismeretek.pptx": "04-uzleti-mukodes",
    "szoftverek_adatbazisok_a_szerzoi_jogi_vedelemben_2026..pptx": "05-szerzoi-jog",
    "Szellemi tulajdonjogok - Iparjogvédelem20260327.pptx": "06-iparjogvedelem",
    "E-kereskedelem_(VIK)kurzus_BVF.pptx": "07-e-kereskedelem",
}


def iter_pictures(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            yield from iter_pictures(sub)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        yield shape


import re

_WS_RE = re.compile(r"\s+")


def _norm(txt):
    return _WS_RE.sub(" ", txt).strip()


_CONT_STARTS = set("()[],;„\"'«»–—…-")


def _is_continuation(txt):
    """Return True if txt looks like a wrapped continuation of the previous paragraph."""
    if not txt:
        return False
    c = txt[0]
    if c.islower():
        return True
    if c in _CONT_STARTS:
        return True
    return False


def shape_items(shape):
    """Yield ordered items: ('p',lvl,txt), ('table',rows), ('img',pic_shape)."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            yield from shape_items(sub)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        yield ("img", shape)
        return
    if shape.has_table:
        tbl = shape.table
        rows = []
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                t = " ".join(_norm(p.text) for p in cell.text_frame.paragraphs if p.text.strip())
                cells.append(t.replace("|", "\\|"))
            rows.append(cells)
        if rows:
            yield ("table", rows)
        return
    if not shape.has_text_frame:
        return

    # First collect paragraphs, then merge wrapped continuations.
    paras = []
    for para in shape.text_frame.paragraphs:
        txt = "".join(run.text for run in para.runs)
        if not txt and para.text:
            txt = para.text
        txt = _norm(txt)
        if not txt:
            continue
        lvl = para.level if para.level is not None else 0
        if paras and paras[-1][0] == lvl and _is_continuation(txt):
            paras[-1] = (lvl, paras[-1][1] + " " + txt)
        else:
            paras.append((lvl, txt))

    for lvl, txt in paras:
        yield ("p", lvl, txt)


def find_logo_hashes(prs, min_count=4, min_ratio=0.35):
    counts = Counter()
    total = 0
    for slide in prs.slides:
        total += 1
        seen = set()
        for shape in slide.shapes:
            for pic in iter_pictures(shape):
                try:
                    seen.add(hashlib.sha256(pic.image.blob).hexdigest())
                except Exception:
                    pass
        for h in seen:
            counts[h] += 1
    return {h for h, c in counts.items() if c >= min_count and c >= total * min_ratio}


def extract_pptx(path, basename):
    prs = Presentation(str(path))
    out_img_dir = IMG_DIR / basename
    out_img_dir.mkdir(parents=True, exist_ok=True)

    logos = find_logo_hashes(prs)
    saved = {}  # hash -> rel path
    img_counter = 0
    skipped_unrenderable = 0

    lines = [f"# {basename}\n"]

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_title = t.splitlines()[0]
                        break

        lines.append(f"\n## Dia {idx}" + (f" — {slide_title}" if slide_title else "") + "\n")

        printed_title = False
        for shape in slide.shapes:
            for item in shape_items(shape):
                kind = item[0]
                if kind == "p":
                    _, lvl, txt = item
                    if not printed_title and slide_title and txt == slide_title:
                        printed_title = True
                        continue
                    lines.append(f"{'  ' * lvl}- {txt}")
                elif kind == "table":
                    rows = item[1]
                    if not rows:
                        continue
                    lines.append("")
                    header = rows[0]
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for r in rows[1:]:
                        r = (r + [""] * len(header))[: len(header)]
                        lines.append("| " + " | ".join(r) + " |")
                    lines.append("")
                elif kind == "img":
                    pic = item[1]
                    try:
                        img = pic.image
                    except Exception:
                        continue
                    blob = img.blob
                    ext = (img.ext or "").lower()
                    h = hashlib.sha256(blob).hexdigest()
                    if h in logos:
                        continue
                    if ext in ("wmf", "emf"):
                        skipped_unrenderable += 1
                        continue
                    if h in saved:
                        rel = saved[h]
                    else:
                        img_counter += 1
                        fname = f"slide-{idx:02d}-{img_counter:03d}.{ext}"
                        (out_img_dir / fname).write_bytes(blob)
                        rel = f"images/{basename}/{fname}"
                        saved[h] = rel
                    lines.append("")
                    lines.append(f"![Dia {idx} kép]({rel})")
                    lines.append("")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append("**Előadói jegyzet:**")
                for ln in notes.splitlines():
                    if ln.strip():
                        lines.append("> " + ln.strip())

    return "\n".join(lines) + "\n", len(saved), skipped_unrenderable


def main():
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    total_imgs = 0
    total_skipped = 0
    for pptx_name, basename in MAPPING.items():
        path = PPT_DIR / pptx_name
        if not path.exists():
            print(f"MISS {pptx_name}", file=sys.stderr)
            continue
        md, n_imgs, n_skip = extract_pptx(path, basename)
        (DOCS_DIR / (basename + ".md")).write_text(md, encoding="utf-8")
        total_imgs += n_imgs
        total_skipped += n_skip
        print(f"OK   {basename:25s} {len(md):6d} chars, {n_imgs:3d} images" +
              (f" ({n_skip} wmf/emf skipped)" if n_skip else ""))
    print(f"\nTotal: {total_imgs} images, {total_skipped} unrenderable skipped")


if __name__ == "__main__":
    main()
